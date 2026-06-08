"""Generate a concise PPT summarizing the empirical analysis.

Outputs: output/empirical_summary.pptx

Slide deck (10 slides):
  1. Title
  2. Problem statement
  3. Pipeline overview
  4. Data foundation (10k → 364 clean fix)
  5. Pattern catalog + match generalization
  6. Backport gap audit
  7. Gap drill: long tail + repo coverage + per-ident
  8. History classification + TRUE backports
  9. Bimodal lag + cross-tab key findings
  10. Caveats + next steps
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path("output") / "empirical_summary.pptx"

# ---- styling -------------------------------------------------------------
ACCENT = RGBColor(0x0E, 0x4F, 0x88)       # deep blue
SUBTLE = RGBColor(0x55, 0x55, 0x55)
NEG = RGBColor(0xC0, 0x39, 0x2B)          # red — caveats / important callouts
POS = RGBColor(0x2E, 0x86, 0x4E)          # green — positive findings


def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, body_lines):
    """body_lines is a list of (text, level, color) tuples."""
    layout = prs.slide_layouts[1]   # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True

    for i, item in enumerate(body_lines):
        text, level, color = item if len(item) == 3 else (item[0], item[1], None)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        for run in p.runs:
            run.font.size = Pt(18 if level == 0 else 16)
            if color:
                run.font.color.rgb = color
    return slide


def add_table_slide(prs, title, headers, rows, note=None):
    layout = prs.slide_layouts[5]   # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    rows_n = len(rows) + 1
    cols_n = len(headers)
    left = Inches(0.5)
    top = Inches(1.4)
    width = Inches(9.0)
    height = Inches(0.35 * rows_n + 0.2)

    tbl = slide.shapes.add_table(rows_n, cols_n, left, top, width, height).table
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(13)
                r.font.color.rgb = ACCENT
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)

    if note:
        txt = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.6), Inches(9.0), Inches(0.6)
        ).text_frame
        txt.text = note
        for p in txt.paragraphs:
            for r in p.runs:
                r.font.size = Pt(12)
                r.font.italic = True
                r.font.color.rgb = SUBTLE
    return slide


# ---- build the deck ------------------------------------------------------


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- 1. Title --------------------------------------------------------
    add_title_slide(
        prs,
        "Workflow Security Fix Backporting",
        "实证分析：模式挖掘 + 跨分支缺口审计\n"
        "样本：Gigawork 10k commits / 364 clean-fix / 359 repos",
    )

    # --- 2. Problem ------------------------------------------------------
    add_content_slide(prs, "研究问题", [
        ("Maintainer 在 master 修了 workflow 安全漏洞", 0, ACCENT),
        ("→ release 分支大概率被遗忘", 1, None),
        ("", 0, None),
        ("目标 1：从历史中挖出可复用的 fix pattern", 0, ACCENT),
        ("→ 给 backport 工具做 catalog", 1, None),
        ("", 0, None),
        ("目标 2：找出 backport 工具未来要去填的洞", 0, ACCENT),
        ("→ 哪些 release 分支当前仍 vulnerable", 1, None),
        ("→ 哪些已经被 backport（ground truth 真值集）", 1, None),
    ])

    # --- 3. Pipeline -----------------------------------------------------
    add_content_slide(prs, "Pipeline 总览（5 阶段，全部确定性）", [
        ("Stage 1  sample       CSV → 10k commits（确定性 hash 采样）", 0, None),
        ("Stage 2  diffs        blob → 14823 YAML 结构化 diff", 0, None),
        ("Stage 3  scan         28k blob × zizmor（in-memory）", 0, None),
        ("Stage 4  clean-fixes  V_fixed≠∅ ∧ V_introduced=∅ → 364", 0, None),
        ("Stage 5  patterns     两级聚类 → 43 桶 / 346 sub-cluster", 0, None),
        ("", 0, None),
        ("Backport 模块（独立目录）", 0, ACCENT),
        ("find-gaps           GitHub API 审 release 分支 → gap/already_fixed/inappl", 1, None),
        ("classify-history    走文件历史确认真 backport + 算 lag", 1, None),
    ])

    # --- 4. Data foundation table ---------------------------------------
    add_table_slide(
        prs, "Stage 1-4：从原始数据到 clean fix",
        ["阶段", "数字", "过滤"],
        [
            ["采样 commit", "10 000", "M change + valid_yaml + valid_workflow"],
            ["file-diff", "14 823", "非空 diff"],
            ["扫描 blob", "28 357", "zizmor JSON 解析成功"],
            ["有任何 fix 的 commit (V_fixed≠∅)", "1 524", "—"],
            ["clean-fix（strict）", "364", "V_introduced == ∅"],
            ["loose-B (备选)", "1 034", "无 ident 净增"],
        ],
        note="脚本：analysis/01_clean_fix_filter_comparison.py"
    )

    # --- 5. Pattern catalog + match -------------------------------------
    add_content_slide(prs, "Stage 5：Pattern catalog + 泛化能力", [
        ("两级聚类（脚本 02）", 0, ACCENT),
        ("L1 主键：frozenset(V_fixed_idents)", 1, None),
        ("L2 主键：结构模板 blake2b hash（保留 [uses=...]）", 1, None),
        ("→ 43 个 level-1 桶 / 346 个 sub-cluster", 1, None),
        ("→ 结构唯一率 0.95（几乎每个 commit 结构独一无二）", 1, NEG),
        ("", 0, None),
        ("Match 实验：fresh 2k commit（seed=99，脚本 03）", 0, ACCENT),
        ("level-1 命中 65/68 = 95.6%（语义类型已饱和）", 1, POS),
        ("level-2 命中 1/68 = 1.5%（结构无法跨 repo 复用）", 1, NEG),
        ("→ Stage 2 必须做元变量参数化", 1, None),
    ])

    # --- 6. Gap audit (KEY FINDING) -------------------------------------
    add_table_slide(
        prs, "Stage 6：Backport 缺口审计（关键发现）",
        ["类别", "数量", "占比"],
        [
            ["release 分支总检查数", "2 546", "100%"],
            ["inapplicable（无文件）", "1 239", "48.7%"],
            ["already_fixed", "472", "18.5%"],
            ["★ gap（仍 vulnerable）", "835", "32.8%"],
            ["—— 可行子集 gap 率", "835 / 1307", "63.9%"],
            ["有 ≥1 gap 的 commit", "101 / 364", "27.7%"],
            ["有 ≥1 gap 的 repo", "98 / 359", "27.3%"],
        ],
        note="脚本：analysis/04_gap_audit_drill.py | 来源：output/backport_gaps/gaps.jsonl"
    )

    # --- 7. Gap drill ---------------------------------------------------
    add_content_slide(prs, "Gap drill：长尾 + repo 分布 + 漏检规则", [
        ("Gap 极端长尾（前 5 个 commit）", 0, ACCENT),
        ("80 gap  archesproject/arches            (unpinned-uses)", 1, None),
        ("44 gap  realm/realm-dotnet              (artipacked+unpinned-uses)", 1, None),
        ("37 gap  datadog/integrations-core       (excessive-permissions)", 1, None),
        ("33 gap  micronaut-projects/micronaut-data (三联)", 1, None),
        ("", 0, None),
        ("Repo 级覆盖（359 audited repos）", 0, ACCENT),
        ("66% 完全无 actionable 信号（无 release 分支或文件不存在）", 1, NEG),
        ("→ backport 只对剩下 34% 的项目有意义", 1, None),
        ("", 0, None),
        ("最常漏 backport 的 zizmor 规则", 0, ACCENT),
        ("unpinned-uses=675   excessive-permissions=316   artipacked=285", 1, None),
    ])

    # --- 8. History + TRUE backport -------------------------------------
    add_table_slide(
        prs, "Stage 7：历史扫描 → 真 backport",
        ["细分状态", "定义", "数量", "占 472"],
        [
            ["★ true_backport", "lag > +1 天", "27", "5.7%"],
            ["same_day_fix", "|lag| ≤ 1 天（多为 merge sync）", "118", "25.0%"],
            ["independent_prior_fix", "lag < -1 天（release 先修）", "6", "1.3%"],
            ["inconclusive", "history cap 不够（MAX=10）", "256", "54.2%"],
            ["never_had_it", "release 历史无此 finding", "17", "3.6%"],
            ["timed_out", "8 分钟硬上限", "48", "10.2%"],
        ],
        note="脚本：analysis/05_history_lag_drill.py"
    )

    # --- 9. Bimodal lag + 27 vs 7 + template-injection ------------------
    add_content_slide(prs, "TRUE backport 的两个反直觉发现", [
        ("Lag 分布双峰（n=27）", 0, ACCENT),
        ("1-7 天：0    1-4 周：0   ← 中间完全空白", 1, NEG),
        ("1-3 月：17（全是 hyperledger/besu 一次 commit）", 1, None),
        ("> 1 年：7（最长 656 天）", 1, None),
        ("→ 真做 backport 的最少也要 ~2 个月", 1, NEG),
        ("", 0, None),
        ("27 vs 7：粒度问题", 0, ACCENT),
        ("27 个 (commit, branch) 对 = backport 工作量", 1, None),
        ("7 个独立 master commit = 独立 fix 事件", 1, None),
        ("→ paper 必须明确选哪个口径报", 1, NEG),
        ("", 0, None),
        ("template-injection 0 真 backport（脚本 06）", 0, ACCENT),
        ("→ release 分支上的 RCE 风险事实上无人维护", 1, NEG),
    ])

    # --- 10. Caveats + next steps ---------------------------------------
    add_content_slide(prs, "局限 + 下一步", [
        ("当前已知局限", 0, ACCENT),
        ("样本只 10k commit（Gigawork 全集 ~1%）", 1, None),
        ("MAX_HISTORY_COMMITS=10 限制了 256 个 inconclusive", 1, None),
        ("same_day_fix 没区分 merge vs cherry-pick vs 真同日修", 1, None),
        ("bot-generated（StepSecurity 等）vs human 没单独 carve out", 1, None),
        ("", 0, None),
        ("下一步建议（按 ROI）", 0, ACCENT),
        ("1. 对 inconclusive 子集用 MAX=50 重跑（30-60 min）", 1, POS),
        ("2. same_day disambiguation（看 commit parents 数）", 1, POS),
        ("3. 扩样本到 50k（一晚上跑完，true_backport → ~135）", 1, POS),
        ("4. 启动 comment 2 baseline（用现有 27 对做 LLM 评测）", 1, POS),
    ])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"slides written -> {OUT}")
    print(f"slide count: {len(prs.slides)}")


if __name__ == "__main__":
    build()
