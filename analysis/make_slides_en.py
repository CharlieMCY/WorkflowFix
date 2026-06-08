"""Generate an English PPT for a non-domain-expert audience.

Outputs: output/empirical_summary_en.pptx

Designed around 5 questions:
  1. How did we select "clean fixes"?
  2. How did we find gap vs truly-backported release branches?
  3. How many repos / clean fixes did we study?
  4. Of clean fixes, how many got backported, how many didn't, how many couldn't be determined?
  5. When backports did happen, how long did they take?

Plain-English vocabulary: no zizmor jargon, no `V_fixed`, no `frozenset`, etc.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

OUT = Path("output") / "empirical_summary_en.pptx"

ACCENT = RGBColor(0x0E, 0x4F, 0x88)
SUBTLE = RGBColor(0x55, 0x55, 0x55)
NEG = RGBColor(0xC0, 0x39, 0x2B)
POS = RGBColor(0x2E, 0x86, 0x4E)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, body):
    """body: list of (text, level, color)."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    for i, item in enumerate(body):
        text, level, color = item if len(item) == 3 else (item[0], item[1], None)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        for r in p.runs:
            r.font.size = Pt(20 if level == 0 else 17)
            if color:
                r.font.color.rgb = color
    return slide


def add_table_slide(prs, title, headers, rows, note=None, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    rows_n = len(rows) + 1
    cols_n = len(headers)
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(0.4 * rows_n + 0.2)

    tbl = slide.shapes.add_table(rows_n, cols_n, left, top, width, height).table
    if col_widths:
        for c, w in enumerate(col_widths):
            tbl.columns[c].width = Inches(w)

    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(14)
                r.font.color.rgb = ACCENT
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(13)

    if note:
        txt = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.6), Inches(9.0), Inches(0.6)
        ).text_frame
        txt.text = note
        for p in txt.paragraphs:
            for r in p.runs:
                r.font.size = Pt(12)
                r.font.italic = True
                r.font.color.rgb = SUBTLE
    return slide


# ---- the deck ------------------------------------------------------------


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(
        prs,
        "Backporting Security Fixes in GitHub Workflows",
        "An empirical study of 359 open-source projects\n"
        "Are old release branches being kept secure?",
    )

    # 2. Background (for non-experts)
    add_content_slide(prs, "Background — in one slide", [
        ("Modern projects use GitHub workflows (YAML files) to automate CI, tests, releases.", 0, None),
        ("These files often have security weaknesses, e.g.:", 0, None),
        ("using third-party actions without pinning them to a specific version", 1, SUBTLE),
        ("giving the build too many permissions", 1, SUBTLE),
        ("letting user-supplied input flow into shell commands", 1, SUBTLE),
        ("", 0, None),
        ("Maintainers fix these on the main branch.", 0, None),
        ("But they often FORGET to apply the same fix to older release branches (e.g. v2.x, v3.x).", 0, NEG),
        ("→ Old versions stay vulnerable.", 0, NEG),
    ])

    # 3. Research questions
    add_content_slide(prs, "Research questions", [
        ("1. Can we automatically identify reliable security fixes in workflow history?", 0, ACCENT),
        ("2. For each such fix on main — do the corresponding release branches get the same fix?", 0, ACCENT),
        ("3. When a backport does happen, how long does it take?", 0, ACCENT),
        ("", 0, None),
        ("Goal: build a dataset of confirmed backport events to support future automation.", 0, None),
    ])

    # 4. Pipeline overview
    add_content_slide(prs, "Pipeline overview", [
        ("Step 1.  Sample 10 000 commits that changed a workflow file (from a public dataset).", 0, None),
        ("Step 2.  Run a static security scanner on the file BEFORE and AFTER each commit.", 0, None),
        ("Step 3.  Identify 'clean fixes' — commits that genuinely removed a security issue.", 0, None),
        ("Step 4.  For each clean fix on the main branch, use the GitHub API to inspect", 0, None),
        ("         every release branch in that project — is the issue still there?", 1, SUBTLE),
        ("Step 5.  Walk each release branch's file history to confirm true backports", 0, None),
        ("         and compute the time gap from the main-branch fix.", 1, SUBTLE),
    ])

    # 5. How we select clean fixes (Q1)
    add_content_slide(prs, "Q1. How is a 'clean fix' defined?", [
        ("For every commit that changed a workflow file:", 0, None),
        ("scan the file BEFORE the commit → list of security issues found", 1, None),
        ("scan the file AFTER the commit → list of remaining issues", 1, None),
        ("", 0, None),
        ("Compute two sets:", 0, ACCENT),
        ("FIXED  =  issues present before, absent after", 1, None),
        ("INTRODUCED  =  issues newly present after", 1, None),
        ("", 0, None),
        ("A commit is a CLEAN FIX iff:", 0, ACCENT),
        ("FIXED is non-empty (the commit really removed something), AND", 1, POS),
        ("INTRODUCED is empty (no new issue accidentally appeared)", 1, POS),
        ("", 0, None),
        ("This conservative criterion filters out false alarms caused by simple reordering / step insertion.", 0, SUBTLE),
    ])

    # 6. Scale (Q3)
    add_table_slide(
        prs, "Q3. Scale of the study",
        ["Stage", "Count"],
        [
            ["Repositories analyzed", "359"],
            ["Commits sampled", "10 000"],
            ["Workflow-file modifications", "14 823"],
            ["File versions scanned by the security tool", "28 357"],
            ["Commits that removed any issue", "1 524"],
            ["Commits passing the clean-fix filter", "364"],
        ],
        col_widths=[6.5, 2.5],
    )

    # 7. How we audit release branches (Q2)
    add_content_slide(prs, "Q2. How we audit release branches", [
        ("For each clean fix on main:", 0, ACCENT),
        ("Use GitHub API to list all release-style branches (release/*, v1.x, stable, …)", 1, None),
        ("Fetch the same workflow file at each release branch's current HEAD", 1, None),
        ("Re-scan it with the security tool", 1, None),
        ("", 0, None),
        ("Classify each release branch into one of three outcomes:", 0, ACCENT),
        ("GAP  →  branch still has the same type of issue (not backported)", 1, NEG),
        ("ALREADY-FIXED  →  branch no longer has it (may or may not be a real backport)", 1, None),
        ("INAPPLICABLE  →  the workflow file doesn't exist on this branch", 1, SUBTLE),
    ])

    # 8. Branch-level results (Q4 part 1)
    add_table_slide(
        prs, "Q4-a. Branch-level results",
        ["Outcome", "Branches", "Share"],
        [
            ["Inapplicable (file absent)", "1 239", "48.7%"],
            ["Already-fixed", "472", "18.5%"],
            ["Still vulnerable (GAP)", "835", "32.8%"],
            ["—", "", ""],
            ["Total release branches audited", "2 546", "100%"],
            ["Among the 'actionable' subset (gap + already-fixed)", "1 307", ""],
            ["→ % of actionable branches that are still vulnerable", "64%", ""],
        ],
        note=(
            "98 of 359 repos (27.3%) have at least one vulnerable release branch."
            "  Long tail: one repo has the same fix missing on 80 release branches."
        ),
        col_widths=[5.5, 2.0, 1.5],
    )

    # 9. Confirming TRUE backports (Q4 part 2 — the why-categorization)
    add_content_slide(prs, "Q4-b. Distinguishing real backports from coincidence", [
        ("'Already-fixed' on a branch ≠ maintainer backported the main-branch fix.", 0, NEG),
        ("It could also mean: that branch never had the issue (different code path).", 0, NEG),
        ("", 0, None),
        ("So we walk each 'already-fixed' branch's file HISTORY:", 0, ACCENT),
        ("Find the commit on the branch where the issue first disappeared.", 1, None),
        ("Compare its date to the main-branch fix date → backport lag.", 1, None),
        ("", 0, None),
        ("Refined categories:", 0, ACCENT),
        ("TRUE BACKPORT  →  branch had the issue, removed it AFTER main fixed (lag > 1 day)", 1, POS),
        ("SAME-DAY  →  removed same day as main (likely auto-merge from main, not a deliberate backport)", 1, SUBTLE),
        ("INDEPENDENT PRIOR FIX  →  branch removed it BEFORE main (release maintainer noticed first)", 1, SUBTLE),
        ("NEVER HAD IT  →  full history shows the issue was never on this branch", 1, SUBTLE),
        ("INCONCLUSIVE / TIMED OUT  →  history too long to determine within budget", 1, SUBTLE),
    ])

    # 10. TRUE backport breakdown (Q4 part 3)
    add_table_slide(
        prs, "Q4-c. Of 472 'already-fixed' branches — what really happened?",
        ["Refined category", "Branches", "Share"],
        [
            ["TRUE BACKPORT (lag > 1 day)", "27", "5.7%"],
            ["Same-day (likely auto-merge)", "118", "25.0%"],
            ["Independent fix on release (lag < -1 day)", "6", "1.3%"],
            ["Inconclusive (file history too deep)", "256", "54.2%"],
            ["Never had the issue", "17", "3.6%"],
            ["Timed out (8-min budget exceeded)", "48", "10.2%"],
        ],
        note=(
            "Important caveat: the 27 TRUE backports come from only ~7 distinct projects. "
            "17 of them are a single fix in hyperledger/besu replayed on 17 release branches."
        ),
        col_widths=[5.5, 2.0, 1.5],
    )

    # 11. Lag distribution (Q5)
    add_table_slide(
        prs, "Q5. How long do backports take? (n = 27)",
        ["Time from main-branch fix to release backport", "Count"],
        [
            ["Within 1 week", "0"],
            ["1 week – 4 weeks", "0"],
            ["1 – 3 months", "17"],
            ["3 – 12 months", "3"],
            ["More than 1 year", "7"],
            ["—", ""],
            ["Median", "51 days"],
            ["Maximum", "656 days"],
        ],
        note=(
            "Key finding: NOBODY backports within a month. When they do, the median is ~2 months. "
            "There is no 'rapid response' cohort."
        ),
        col_widths=[7.5, 1.5],
    )

    # 12. Summary
    add_content_slide(prs, "Summary", [
        ("Studied 359 open-source projects / 364 verified clean security fixes.", 0, None),
        ("", 0, None),
        ("Key empirical findings:", 0, ACCENT),
        ("27% of projects have at least one release branch still vulnerable", 1, NEG),
        ("Of every audited (fix, release branch) pair, 33% are unpatched", 1, NEG),
        ("Only 27 backport events confirmed — and from only ~7 distinct projects", 1, NEG),
        ("Backports never happen within a month; median delay is ~2 months", 1, NEG),
        ("", 0, None),
        ("Implication: an automated backporting tool would have abundant real-world demand,", 0, ACCENT),
        ("                 and there is ample low-hanging fruit on existing release branches.", 0, ACCENT),
    ])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"slides written -> {OUT}")
    print(f"slide count:    {len(prs.slides)}")


if __name__ == "__main__":
    build()
